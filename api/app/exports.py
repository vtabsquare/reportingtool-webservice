from __future__ import annotations
from pathlib import Path
from io import BytesIO
import json, smtplib, os, urllib.request, urllib.error
from email.message import EmailMessage

# pptx and reportlab are imported lazily (inside each function) so that the
# backend server starts successfully even if these optional packages are not
# bundled by PyInstaller. Core login, data, and report-authoring features are
# completely unaffected. Only PDF/PPT export calls would raise an error if the
# package is genuinely absent from the bundle.


def _pages(project):
    return (project or {}).get('report',{}).get('pages',[]) or []

def report_pdf(project:dict)->bytes:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.pdfgen import canvas
    except ImportError as e:
        raise RuntimeError(f'PDF export is unavailable: reportlab is not installed in the desktop bundle. ({e})') from e
    buf=BytesIO(); c=canvas.Canvas(buf,pagesize=landscape(A4)); W,H=landscape(A4)
    name=project.get('report',{}).get('name') or project.get('name') or 'VTAB Report'
    pages=_pages(project) or [{'name':'Page 1','visuals':[]}]
    for i,p in enumerate(pages,1):
        c.setFillColor(colors.HexColor('#081525'));c.rect(0,H-42,W,42,stroke=0,fill=1)
        c.setFillColor(colors.white);c.setFont('Helvetica-Bold',18);c.drawString(28,H-28,name)
        c.setFillColor(colors.HexColor('#111827'));c.setFont('Helvetica-Bold',15);c.drawString(28,H-68,p.get('name') or f'Page {i}')
        y=H-100
        for v in (p.get('visuals') or [])[:18]:
            title=v.get('title') or v.get('type','Visual').replace('_',' ').title(); typ=v.get('type','visual')
            c.setFillColor(colors.HexColor('#f8fafc'));c.roundRect(28,y-38,W-56,32,8,stroke=1,fill=1)
            c.setFillColor(colors.HexColor('#0f172a'));c.setFont('Helvetica-Bold',10);c.drawString(40,y-20,title)
            c.setFillColor(colors.HexColor('#64748b'));c.setFont('Helvetica',8);c.drawRightString(W-40,y-20,typ)
            y-=42
            if y<40: break
        c.setFillColor(colors.HexColor('#64748b'));c.setFont('Helvetica',8);c.drawRightString(W-28,18,f'Page {i} of {len(pages)} · Exported from VTAB Reporting Studio')
        c.showPage()
    c.save();return buf.getvalue()

def report_pptx(project:dict)->bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise RuntimeError(f'PowerPoint export is unavailable: python-pptx is not installed in the desktop bundle. ({e})') from e
    prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)
    name=project.get('report',{}).get('name') or project.get('name') or 'VTAB Report'
    pages=_pages(project) or [{'name':'Page 1','visuals':[]}]
    for p in pages:
        slide=prs.slides.add_slide(prs.slide_layouts[6])
        title=slide.shapes.add_textbox(Inches(.5),Inches(.3),Inches(12.3),Inches(.6));tf=title.text_frame;tf.text=name;tf.paragraphs[0].font.size=Pt(24);tf.paragraphs[0].font.bold=True
        sub=slide.shapes.add_textbox(Inches(.5),Inches(.95),Inches(12.3),Inches(.35));st=sub.text_frame;st.text=p.get('name','Report Page');st.paragraphs[0].font.size=Pt(13)
        visuals=p.get('visuals') or []; cols=2; card_w=6.05; card_h=1.1; x0=.5;y0=1.45
        for idx,v in enumerate(visuals[:10]):
            row=idx//cols;col=idx%cols;x=x0+col*6.2;y=y0+row*1.18
            shp=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(card_w),Inches(card_h));tf=shp.text_frame
            tf.text=v.get('title') or v.get('type','Visual').replace('_',' ').title();tf.paragraphs[0].font.bold=True;tf.paragraphs[0].font.size=Pt(12)
            p2=tf.add_paragraph();p2.text=f"Type: {v.get('type','visual')}";p2.font.size=Pt(9)
            vals=v.get('measures') or v.get('values') or []
            if vals:
                p3=tf.add_paragraph();p3.text='Measures: '+', '.join(map(str,vals[:4]));p3.font.size=Pt(8)
    buf=BytesIO();prs.save(buf);return buf.getvalue()

def _send_via_brevo(*, api_key:str, from_email:str, to:list[str], subject:str, body:str):
    """Send a plain-text email using the Brevo (Sendinblue) Transactional Email REST API."""
    payload = json.dumps({
        'sender': {'email': from_email},
        'to': [{'email': addr} for addr in to],
        'subject': subject,
        'textContent': body,
    }).encode('utf-8')
    req = urllib.request.Request(
        'https://api.brevo.com/v3/smtp/email',
        data=payload,
        method='POST',
        headers={
            'api-key': api_key,
            'Content-Type': 'application/json',
            'Accept': 'application/json',
        }
    )
    try:
        with urllib.request.urlopen(req, timeout=20) as r:
            r.read()
    except urllib.error.HTTPError as e:
        err_body = e.read().decode('utf-8', errors='replace')
        raise ValueError(f'Brevo API error {e.code}: {err_body}')
    return {'ok': True, 'recipients': to}

def send_report_email(*,smtp_cfg:dict,to:list[str],subject:str,body:str,report_url:str|None=None,attachment_name:str|None=None,attachment:bytes|None=None,attachment_type:str|None=None):
    content=body.strip() or 'A VTAB report has been shared with you.'
    if report_url:content+=f"\n\nOpen report: {report_url}"

    # Use Brevo REST API if BREVO_API_KEY is configured (no SMTP needed)
    brevo_key = os.environ.get('BREVO_API_KEY') or smtp_cfg.get('brevoApiKey')
    if brevo_key:
        from_email = smtp_cfg.get('fromEmail') or smtp_cfg.get('username') or os.environ.get('VTAB_SMTP_FROM','')
        if not from_email: raise ValueError('Set VTAB_SMTP_FROM (the verified sender email) in your .env file.')
        # Attachments via Brevo API require base64 encoding — for now fall through to SMTP if attachment present
        if not attachment:
            return _send_via_brevo(api_key=brevo_key, from_email=from_email, to=to, subject=subject, body=content)

    # Fall back to SMTP (also handles attachment emails)
    if not smtp_cfg.get('host'): raise ValueError('SMTP host is not configured. Add BREVO_API_KEY or configure SMTP in Admin > Email & Workspace Settings.')
    msg=EmailMessage();msg['Subject']=subject;msg['From']=smtp_cfg.get('fromEmail') or smtp_cfg.get('username');msg['To']=', '.join(to)
    msg.set_content(content)
    if attachment and attachment_name:
        maintype,subtype=('application','octet-stream')
        if attachment_type=='pdf':subtype='pdf'
        elif attachment_type=='pptx':subtype='vnd.openxmlformats-officedocument.presentationml.presentation'
        msg.add_attachment(attachment,maintype=maintype,subtype=subtype,filename=attachment_name)
    port=int(smtp_cfg.get('port') or (465 if smtp_cfg.get('ssl') else 587))
    cls=smtplib.SMTP_SSL if smtp_cfg.get('ssl') else smtplib.SMTP
    with cls(smtp_cfg['host'],port,timeout=20) as s:
        if not smtp_cfg.get('ssl') and smtp_cfg.get('startTls',True):s.starttls()
        if smtp_cfg.get('username'):s.login(smtp_cfg['username'],smtp_cfg.get('password',''))
        s.send_message(msg)
    return {'ok':True,'recipients':to}
